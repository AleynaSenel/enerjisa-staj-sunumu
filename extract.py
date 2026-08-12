from pptx import Presentation
prs = Presentation(r"c:\Users\Aleyna\Desktop\toroslar-tanitim\Enerjisa_Sunum spark 2026 (2).pptx")
with open('ppt_content.txt', 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"\n--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                f.write(shape.text + "\n")
